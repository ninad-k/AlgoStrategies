"""Build the Automation Department presentation with Rey Capital branding.

Run:
    python3 build_deck.py

Produces Automation_Department_Overview.pptx alongside this script.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---- Brand ---------------------------------------------------------------
BRAND_BLUE = RGBColor(0x00, 0x4A, 0xAC)
BRAND_LIGHT = RGBColor(0xBD, 0xD4, 0xF5)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MUTED_TEXT = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BG = RGBColor(0xF4, 0xF7, 0xFC)

FONT = "Calibri"

# ---- Paths ---------------------------------------------------------------
HERE = Path(__file__).resolve().parent
REPO_ROOT = HERE.parents[2]
ASSETS = REPO_ROOT / "monitoring" / "dashboards" / "mt5_pnl_dashboard" / "src" / "TradeAtlas" / "Assets"
LOGO_COLOR = ASSETS / "ReyCapital_Logo.png"
LOGO_WHITE = ASSETS / "ReyCapital_Logo_White.png"
OUTPUT = HERE / "Automation_Department_Overview.pptx"

# ---- Layout constants (16:9) ---------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(1.1)
FOOTER_H = Inches(0.35)


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_fill(shape, color)
    return shape


def _add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = DARK_TEXT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_bullets(slide, bullets, left, top, width, height, *, size: int = 18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = DARK_TEXT
    return box


def _add_header(slide, title: str) -> None:
    # Brand blue header band.
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, HEADER_H, BRAND_BLUE)
    # White mini logo top-left.
    logo_h = Inches(0.65)
    slide.shapes.add_picture(
        str(LOGO_WHITE),
        Inches(0.45),
        Inches(0.225),
        height=logo_h,
    )
    # Title text, right side of the band.
    _add_text(
        slide,
        title,
        left=Inches(3.6),
        top=Inches(0.25),
        width=Inches(9.4),
        height=Inches(0.7),
        size=30,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Thin accent underline.
    _add_rect(
        slide,
        Emu(0),
        HEADER_H,
        SLIDE_W,
        Inches(0.06),
        BRAND_LIGHT,
    )


def _add_footer(slide, page_no: int, total: int) -> None:
    _add_rect(
        slide,
        Emu(0),
        SLIDE_H - FOOTER_H,
        SLIDE_W,
        FOOTER_H,
        BRAND_BLUE,
    )
    _add_text(
        slide,
        "Rey Capital  •  Smart Investments",
        left=Inches(0.4),
        top=SLIDE_H - FOOTER_H,
        width=Inches(8),
        height=FOOTER_H,
        size=10,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_text(
        slide,
        f"{page_no} / {total}",
        left=SLIDE_W - Inches(1.2),
        top=SLIDE_H - FOOTER_H,
        width=Inches(0.8),
        height=FOOTER_H,
        size=10,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ---- Slide builders ------------------------------------------------------

def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Soft background.
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, SOFT_BG)
    # Top blue band.
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.6), BRAND_BLUE)
    # Bottom blue band.
    _add_rect(slide, Emu(0), SLIDE_H - Inches(0.6), SLIDE_W, Inches(0.6), BRAND_BLUE)

    # Centered color logo.
    logo_w = Inches(4.5)
    slide.shapes.add_picture(
        str(LOGO_COLOR),
        (SLIDE_W - logo_w) / 2,
        Inches(1.4),
        width=logo_w,
    )

    _add_text(
        slide,
        "Automation Department",
        left=Inches(1),
        top=Inches(3.9),
        width=Inches(11.333),
        height=Inches(1),
        size=48,
        bold=True,
        color=BRAND_BLUE,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Overview & Roadmap",
        left=Inches(1),
        top=Inches(4.95),
        width=Inches(11.333),
        height=Inches(0.6),
        size=24,
        color=MUTED_TEXT,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Presented by  <Your Name>",
        left=Inches(1),
        top=Inches(5.8),
        width=Inches(11.333),
        height=Inches(0.5),
        size=18,
        color=DARK_TEXT,
        align=PP_ALIGN.CENTER,
    )


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    page_no: int,
    total: int,
    *,
    subtitle: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, WHITE)
    _add_header(slide, title)

    top = Inches(1.55)
    if subtitle:
        _add_text(
            slide,
            subtitle,
            left=Inches(0.7),
            top=top,
            width=Inches(12),
            height=Inches(0.5),
            size=18,
            color=MUTED_TEXT,
        )
        top = Inches(2.1)

    _add_bullets(
        slide,
        bullets,
        left=Inches(0.8),
        top=top,
        width=Inches(11.7),
        height=Inches(4.8),
        size=22,
    )
    _add_footer(slide, page_no, total)


def add_closing_slide(prs: Presentation, page_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, BRAND_BLUE)

    logo_w = Inches(4.2)
    slide.shapes.add_picture(
        str(LOGO_WHITE),
        (SLIDE_W - logo_w) / 2,
        Inches(1.2),
        width=logo_w,
    )

    _add_text(
        slide,
        "Thank You",
        left=Inches(1),
        top=Inches(3.7),
        width=Inches(11.333),
        height=Inches(1.2),
        size=60,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Q & A   •   Live Demo   •   Open Floor",
        left=Inches(1),
        top=Inches(5.0),
        width=Inches(11.333),
        height=Inches(0.6),
        size=22,
        color=BRAND_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        f"{page_no} / {total}",
        left=SLIDE_W - Inches(1.2),
        top=SLIDE_H - Inches(0.5),
        width=Inches(0.8),
        height=Inches(0.35),
        size=10,
        color=BRAND_LIGHT,
        align=PP_ALIGN.RIGHT,
    )


# ---- Deck content --------------------------------------------------------

CONTENT_SLIDES: list[tuple[str, str | None, list[str]]] = [
    (
        "Agenda",
        "What we'll walk through today",
        [
            "Department Overview",
            "Tools & Technologies",
            "Achievements & Metrics",
            "Current Projects",
            "Challenges & Lessons Learned",
            "Future Roadmap",
            "Trading / Algo-Specific Automation",
            "Q&A / Demo",
        ],
    ),
    (
        "Department Overview",
        "Who we are and what we do",
        [
            "Mission, vision, and role within the organization",
            "Team structure and key members",
            "Current projects and initiatives",
        ],
    ),
    (
        "Tools & Technologies",
        "Our modern automation stack",
        [
            "AI / ML integration in automation",
            "Leveraging LLMs for code generation, review, and ops",
            "Intelligent decisioning inside automated workflows",
            "Continuous experimentation and model feedback loops",
        ],
    ),
    (
        "Achievements & Metrics",
        "Measurable impact delivered",
        [
            "Hours saved  /  ROI delivered",
            "Error reduction rates",
            "Key success stories  /  case studies",
            "Before vs. after comparisons",
        ],
    ),
    (
        "Current Projects",
        "What we're actively building",
        [
            "In-flight automations",
            "Pipeline of upcoming work",
            "Cross-team collaborations",
        ],
    ),
    (
        "Challenges & Lessons Learned",
        "What we're solving for",
        [
            "Common bottlenecks",
            "Change management",
            "Maintenance overhead",
        ],
    ),
    (
        "Future Roadmap",
        "Where we're headed next",
        [
            "AI-driven  /  intelligent automation",
            "Hyperautomation strategy",
            "Upskilling and training plans",
        ],
    ),
    (
        "Trading / Algo-Specific Automation",
        "Automation applied to our trading stack",
        [
            "EA  /  strategy automation",
            "Backtesting pipelines",
            "Broker connector automation",
            "Monitoring dashboards",
        ],
    ),
    (
        "Q&A  /  Demo",
        "Let's see it in action",
        [
            "Live demo of a flagship automation",
            "Open floor for questions and discussion",
        ],
    ),
]


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 1 + len(CONTENT_SLIDES) + 1  # title + content + closing

    add_title_slide(prs)
    for idx, (title, subtitle, bullets) in enumerate(CONTENT_SLIDES, start=2):
        add_content_slide(prs, title, bullets, idx, total, subtitle=subtitle)
    add_closing_slide(prs, total, total)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
